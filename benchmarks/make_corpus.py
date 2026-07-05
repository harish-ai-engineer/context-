"""Generate the golden corpus: source documents + golden JSON, from one definition.

Each corpus entry is defined once as structured content and rendered to its
source format (md/html/txt/pdf/docx) AND to a golden file — so the golden is
correct by construction and the benchmark measures parse-back fidelity.

Usage:  python benchmarks/make_corpus.py
"""

from __future__ import annotations

import json
import os

HERE = os.path.dirname(os.path.abspath(__file__))
CORPUS = os.path.join(HERE, "corpus")
GOLDEN = os.path.join(CORPUS, "golden")

# ---------------------------------------------------------------------------
# Content definitions: (kind, *args) items in reading order.
#   ("h", level, title)        heading
#   ("p", text)                paragraph
#   ("li", text)               list item
#   ("q", text)                quote
# ---------------------------------------------------------------------------

QUARTERLY_REPORT = [
    ("h", 1, "Q2 2026 Quarterly Report"),
    ("p", "This report summarizes Aurora Systems' performance for the second "
          "quarter of fiscal year 2026, covering April through June."),
    ("h", 2, "Financial Highlights"),
    ("p", "Revenue grew 12 percent year over year to 48.3 million USD, driven "
          "primarily by expansion in the enterprise segment."),
    ("p", "Gross margin improved to 71 percent, up from 68 percent in the "
          "prior quarter, reflecting infrastructure cost reductions."),
    ("h", 3, "Revenue by Segment"),
    ("p", "Enterprise contracts contributed 29.1 million USD while "
          "self-service subscriptions accounted for 19.2 million USD."),
    ("h", 2, "Operational Update"),
    ("p", "Headcount reached 412 employees at quarter end, with 38 net new "
          "hires concentrated in engineering and customer success."),
    ("li", "Launched the Aurora Edge deployment option in May."),
    ("li", "Achieved SOC 2 Type II certification in June."),
    ("li", "Opened the Singapore regional office."),
    ("h", 2, "Outlook"),
    ("p", "Management expects third quarter revenue between 50 and 52 million "
          "USD, assuming stable enterprise renewal rates."),
]

REFUND_POLICY = [
    ("h", 1, "Refund Policy"),
    ("p", "Customers may request a full refund within 30 days of purchase. "
          "Refunds are processed to the original payment method within five "
          "business days of approval."),
    ("h", 2, "Eligibility"),
    ("p", "To be eligible, the product must be unused and the request must "
          "include the original order number."),
    ("li", "Annual plans: prorated refund after 30 days."),
    ("li", "Monthly plans: no partial-month refunds."),
    ("h", 2, "Exceptions"),
    ("p", "Digital goods marked as final sale are excluded from this policy, "
          "as are purchases made through third-party resellers."),
]

API_DOCS = [
    ("h", 1, "Ingest API Reference"),
    ("p", "The Ingest API accepts documents for parsing and returns a job "
          "identifier that can be polled for results."),
    ("h", 2, "Authentication"),
    ("p", "All requests require a bearer token passed in the Authorization "
          "header. Tokens are scoped per project and expire after 90 days."),
    ("h", 2, "Endpoints"),
    ("h", 3, "POST /v1/parse"),
    ("p", "Submits a document for parsing. The request body is multipart form "
          "data with a single file field."),
    ("li", "file: the document to parse, up to 50 MB."),
    ("li", "callback_url: optional webhook for completion."),
    ("h", 3, "GET /v1/jobs"),
    ("p", "Lists recent parse jobs for the authenticated project, newest "
          "first, paginated in groups of 25."),
]

NEWS_ARTICLE = [
    ("h", 1, "City Council Approves Riverfront Renewal Plan"),
    ("p", "After eight months of public hearings, the city council voted 7 to "
          "2 on Tuesday to approve the riverfront renewal plan."),
    ("p", "The 140 million dollar project will convert the disused freight "
          "yard into a park, housing, and a flood-control wetland."),
    ("q", "This is the largest green investment in the city's history, and it "
          "will pay for itself in avoided flood damage alone."),
    ("p", "Construction is expected to begin next spring, with the first "
          "phase of the park opening two years later."),
]

MEETING_NOTES = [
    ("p", "Weekly platform sync, July 1st. Attendees: Dana, Marcus, Priya, "
          "and Tom. Dana chaired the meeting."),
    ("p", "Priya reported that the ingestion backlog dropped from 4,200 to "
          "300 documents after the parser fix shipped on Friday."),
    ("p", "Marcus raised a concern about provenance loss in the HTML "
          "pipeline; a spike was scheduled for next sprint."),
    ("p", "Action items: Tom to draft the OCR vendor comparison, Dana to "
          "schedule the security review, Priya to close out the backlog "
          "dashboard."),
]

INVOICE_PAGES = [  # one list of lines per PDF page
    [
        "INVOICE 2041",
        "Billed to Acme Corporation for consulting services rendered in June 2026.",
        "Engagement: data pipeline architecture review, 32 hours at 120 USD per hour.",
        "Subtotal 3,840 USD. Tax at 8 percent adds 307 USD.",
    ],
    [
        "Total amount due: 4,147 USD within 30 days of the invoice date.",
        "Payment accepted by bank transfer to the account listed in the master services agreement.",
        "Late payments accrue interest at 1.5 percent per month as permitted by law.",
    ],
]

HR_HANDBOOK = [
    ("h", 1, "Employee Handbook"),
    ("p", "This handbook summarizes working policies for all full-time "
          "employees and applies from the first day of employment."),
    ("h", 2, "Time Off"),
    ("p", "Employees accrue 1.75 vacation days per month, capped at 30 days. "
          "Unused days above the cap are forfeited at year end."),
    ("h", 2, "Remote Work"),
    ("p", "Employees may work remotely up to three days per week with "
          "manager approval recorded in the HR system."),
]
HR_TABLE = [
    ["Benefit", "Eligibility", "Waiting period"],
    ["Health insurance", "All employees", "None"],
    ["401k matching", "Full-time", "90 days"],
    ["Sabbatical", "5+ years tenure", "None"],
]


# ------------------------------------------------------------- renderers --

def _golden(items, *, pages=None, tables=None, prov_samples=None) -> dict:
    texts = []
    headings = []
    for it in items:
        if it[0] == "h":
            headings.append([it[1], it[2]])
            texts.append(it[2])
        else:
            texts.append(it[1])
    return {
        "text": " ".join(" ".join(texts).split()),
        "headings": headings,
        "tables": tables or [],
        "provenance": prov_samples or [],
        "pages": pages,
    }


def render_md(items) -> str:
    out = []
    for it in items:
        if it[0] == "h":
            out.append("#" * it[1] + " " + it[2])
        elif it[0] == "li":
            out.append("- " + it[1])
        elif it[0] == "q":
            out.append("> " + it[1])
        else:
            out.append(it[1])
    return "\n\n".join(out) + "\n"


def render_html(items, title: str) -> str:
    out = [f"<html><head><title>{title}</title></head><body>"]
    open_list = False
    for it in items:
        if it[0] == "li" and not open_list:
            out.append("<ul>")
            open_list = True
        if it[0] != "li" and open_list:
            out.append("</ul>")
            open_list = False
        if it[0] == "h":
            out.append(f"<h{it[1]}>{it[2]}</h{it[1]}>")
        elif it[0] == "li":
            out.append(f"<li>{it[1]}</li>")
        elif it[0] == "q":
            out.append(f"<blockquote>{it[1]}</blockquote>")
        else:
            out.append(f"<p>{it[1]}</p>")
    if open_list:
        out.append("</ul>")
    out.append("</body></html>")
    return "\n".join(out) + "\n"


def render_txt(items) -> str:
    return "\n\n".join(it[1] for it in items if it[0] == "p") + "\n"


def render_pdf(pages: list[list[str]]) -> bytes:
    """Hand-rolled minimal PDF: one content stream of Tj lines per page."""
    objects: list[bytes] = []

    def esc(s: str) -> str:
        return s.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")

    n_pages = len(pages)
    font_obj = 3 + 2 * n_pages
    kids = " ".join(f"{3 + 2 * i} 0 R" for i in range(n_pages))
    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>")
    objects.append(f"<< /Type /Pages /Kids [{kids}] /Count {n_pages} >>".encode())
    for i, lines in enumerate(pages):
        page_no = 3 + 2 * i
        objects.append(
            (f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
             f"/Resources << /Font << /F1 {font_obj} 0 R >> >> "
             f"/Contents {page_no + 1} 0 R >>").encode()
        )
        parts = ["BT /F1 11 Tf 72 720 Td"]
        for j, line in enumerate(lines):
            if j:
                parts.append("0 -18 Td")
            parts.append(f"({esc(line)}) Tj")
        parts.append("ET")
        stream = " ".join(parts).encode()
        objects.append(b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n"
                       + stream + b"\nendstream")
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    buf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for i, body in enumerate(objects, 1):
        offsets.append(len(buf))
        buf += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_at = len(buf)
    buf += f"xref\n0 {len(objects) + 1}\n".encode()
    buf += b"0000000000 65535 f \n"
    for off in offsets[1:]:
        buf += f"{off:010d} 00000 n \n".encode()
    buf += (f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_at}\n%%EOF\n").encode()
    return bytes(buf)


def render_pptx(path: str, slides: list[list[str]]) -> None:
    """Valid PPTX via python-pptx so every tool under benchmark can open it."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for texts in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = texts[0]
        body = slide.placeholders[1].text_frame
        for j, para in enumerate(texts[1:]):
            (body.paragraphs[0] if j == 0 else body.add_paragraph()).text = para
    prs.save(path)


def render_xlsx(path: str, sheet_name: str, rows: list[list[str]]) -> None:
    """Valid XLSX via openpyxl (exercises the sharedStrings path in our parser)."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    for row in rows:
        ws.append(row)
    wb.save(path)


PRODUCT_DECK = [
    ["Atlas Data Platform", "A unified layer for ingesting, cataloging, and governing analytical data."],
    ["Lineage Tracking", "Column-level lineage is captured automatically for every pipeline run.",
     "Auditors can trace any dashboard number back to its source system."],
    ["Pricing", "Standard tier starts at 2,000 USD per month with unlimited seats."],
]

SALES_ROWS = [
    ["Region", "Units", "Revenue"],
    ["North", "1240", "310000"],
    ["South", "980", "245000"],
    ["Europe", "1610", "402500"],
]


def render_docx(path: str, items, table_rows) -> None:
    import docx

    d = docx.Document()
    for it in items:
        if it[0] == "h":
            d.add_heading(it[2], level=it[1])
        else:
            d.add_paragraph(it[1])
    t = d.add_table(rows=len(table_rows), cols=len(table_rows[0]))
    for r, row in enumerate(table_rows):
        for c, val in enumerate(row):
            t.cell(r, c).text = val
    d.save(path)


# ------------------------------------------------------------------ main --

def main() -> None:
    os.makedirs(GOLDEN, exist_ok=True)

    def write(name: str, data, golden: dict) -> None:
        path = os.path.join(CORPUS, name)
        if isinstance(data, bytes):
            with open(path, "wb") as fh:
                fh.write(data)
        elif data is not None:
            with open(path, "w", encoding="utf-8", newline="\n") as fh:
                fh.write(data)
        stem = os.path.splitext(name)[0]
        with open(os.path.join(GOLDEN, stem + ".json"), "w", encoding="utf-8") as fh:
            json.dump(golden, fh, indent=2, ensure_ascii=False)
        print("wrote", name)

    write("quarterly-report.md", render_md(QUARTERLY_REPORT), _golden(
        QUARTERLY_REPORT,
        prov_samples=[
            {"prefix": "Enterprise contracts contributed",
             "section_path": "Q2 2026 Quarterly Report > Financial Highlights > Revenue by Segment"},
            {"prefix": "Management expects third quarter",
             "section_path": "Q2 2026 Quarterly Report > Outlook"},
        ]))

    write("refund-policy.md", render_md(REFUND_POLICY), _golden(
        REFUND_POLICY,
        prov_samples=[
            {"prefix": "Digital goods marked as final sale",
             "section_path": "Refund Policy > Exceptions"},
        ]))

    write("api-docs.html", render_html(API_DOCS, "Ingest API Reference"), _golden(
        API_DOCS,
        prov_samples=[
            {"prefix": "Submits a document for parsing",
             "section_path": "Ingest API Reference > Endpoints > POST /v1/parse"},
            {"prefix": "Lists recent parse jobs",
             "section_path": "Ingest API Reference > Endpoints > GET /v1/jobs"},
        ]))

    write("news-article.html", render_html(NEWS_ARTICLE, "Riverfront Renewal"), _golden(
        NEWS_ARTICLE,
        prov_samples=[
            {"prefix": "Construction is expected to begin",
             "section_path": "City Council Approves Riverfront Renewal Plan"},
        ]))

    write("meeting-notes.txt", render_txt(MEETING_NOTES), _golden(
        MEETING_NOTES,
        prov_samples=[
            {"prefix": "Weekly platform sync", "section_path": None},
        ]))

    pdf_items = [("p", line) for page in INVOICE_PAGES for line in page]
    write("invoice.pdf", render_pdf(INVOICE_PAGES), _golden(
        pdf_items, pages=len(INVOICE_PAGES),
        prov_samples=[
            {"prefix": "INVOICE 2041", "page": 1},
            {"prefix": "Total amount due", "page": 2},
        ]))

    # Slide-title heading LEVEL is a rendering convention (we use h2, others
    # use h1), so decks don't participate in the structure metric — slide
    # mapping is scored by the provenance samples instead.
    deck_items = [("p", t) for texts in PRODUCT_DECK for t in texts]
    try:
        render_pptx(os.path.join(CORPUS, "product-deck.pptx"), PRODUCT_DECK)
        write("product-deck.pptx", None, _golden(
            deck_items, pages=len(PRODUCT_DECK),
            prov_samples=[
                {"prefix": "Column-level lineage", "page": 2, "section_path": "Lineage Tracking"},
                {"prefix": "Standard tier starts", "page": 3, "section_path": "Pricing"},
            ]))
    except ImportError:
        print("skipped product-deck.pptx (python-pptx not installed)")

    try:
        render_xlsx(os.path.join(CORPUS, "sales-summary.xlsx"), "Q3 Sales", SALES_ROWS)
        write("sales-summary.xlsx", None, _golden(
            [], tables=[SALES_ROWS],
            prov_samples=[
                {"prefix": "| Region", "section_path": "Q3 Sales"},
            ]))
    except ImportError:
        print("skipped sales-summary.xlsx (openpyxl not installed)")

    try:
        render_docx(os.path.join(CORPUS, "hr-handbook.docx"), HR_HANDBOOK, HR_TABLE)
        write("hr-handbook.docx", None, _golden(
            HR_HANDBOOK, tables=[HR_TABLE],
            prov_samples=[
                {"prefix": "Employees may work remotely",
                 "section_path": "Employee Handbook > Remote Work"},
            ]))
    except ImportError:
        print("skipped hr-handbook.docx (python-docx not installed)")


if __name__ == "__main__":
    main()
